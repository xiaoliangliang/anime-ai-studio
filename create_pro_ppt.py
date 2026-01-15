# -*- coding: utf-8 -*-
"""
DramaAI 世界级黑客松PPT生成器
使用 python-pptx 创建专业PowerPoint演示文稿
设计理念: Data Flow Chromatics
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os

# 输出文件
OUTPUT_FILE = "d:/projects/dramaai/DramaAI_世界级路演PPT.pptx"

# Data Flow Chromatics 色彩系统
class Colors:
    BACKGROUND = RGBColor(10, 14, 39)  # #0a0e27
    TEXT_PRIMARY = RGBColor(255, 255, 255)
    TEXT_SECONDARY = RGBColor(160, 174, 192)

    # 五阶段渐变色（起始色）
    SCREENWRITER = RGBColor(102, 126, 234)  # #667eea
    STORYBOARD = RGBColor(240, 147, 251)    # #f093fb
    DESIGNER = RGBColor(245, 87, 108)       # #f5576c
    ARTIST = RGBColor(79, 172, 254)         # #4facfe
    DIRECTOR = RGBColor(0, 242, 254)        # #00f2fe

    ACCENT = RGBColor(0, 242, 254)

def set_slide_background(slide, color):
    """设置幻灯片背景色"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title(slide, text, font_size=36, bold=True, color=None):
    """添加标题"""
    if color is None:
        color = Colors.TEXT_PRIMARY

    left = Inches(1)
    top = Inches(0.5)
    width = Inches(8)
    height = Inches(1)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.text = text
    text_frame.word_wrap = True

    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.font.size = Pt(font_size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    paragraph.font.name = "微软雅黑"

    return textbox

def add_centered_text(slide, text, top_inches, font_size=24, color=None):
    """添加居中文字"""
    if color is None:
        color = Colors.TEXT_PRIMARY

    textbox = slide.shapes.add_textbox(Inches(1), top_inches, Inches(8), Inches(0.5))
    text_frame = textbox.text_frame
    text_frame.text = text

    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.font.size = Pt(font_size)
    paragraph.font.color.rgb = color
    paragraph.font.name = "微软雅黑"

    return textbox

def create_gradient_box(slide, left, top, width, height, start_color, end_color):
    """创建渐变效果矩形（模拟）"""
    shape = slide.shapes.add_shape(
        1,  # Rectangle
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = start_color
    shape.line.color.rgb = start_color

    return shape

def create_cover_slide(prs):
    """第1页：封面"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    set_slide_background(slide, Colors.BACKGROUND)

    # 装饰圆（右侧）
    shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(6), Inches(2), Inches(3), Inches(3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = Colors.SCREENWRITER
    shape.line.fill.background()

    # 主标题
    add_title(slide, "DramaAI", font_size=60, color=Colors.TEXT_PRIMARY)

    # 副标题
    add_centered_text(slide, "AI短剧一站式创作平台", Inches(2.2),
                     font_size=28, color=Colors.ACCENT)

    # 五阶段标签
    stages = ["编剧", "分镜", "设计", "美工", "导演"]
    stage_colors = [Colors.SCREENWRITER, Colors.STORYBOARD, Colors.DESIGNER,
                    Colors.ARTIST, Colors.DIRECTOR]

    start_x = Inches(0.8)
    y = Inches(4.5)
    box_size = Inches(0.8)
    spacing = Inches(1.5)

    for i, (stage, color) in enumerate(zip(stages, stage_colors)):
        x = start_x + i * spacing
        box = create_gradient_box(slide, x, y, box_size, box_size, color, color)

        # 阶段名称
        textbox = slide.shapes.add_textbox(x, y + box_size + Inches(0.1),
                                           box_size, Inches(0.3))
        tf = textbox.text_frame
        tf.text = stage
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].font.color.rgb = Colors.TEXT_PRIMARY
        tf.paragraphs[0].font.name = "微软雅黑"

    # Slogan
    add_centered_text(slide, "让每个人都能成为动漫导演", Inches(6),
                     font_size=18, color=Colors.TEXT_SECONDARY)

def create_problem_slide(prs):
    """第2页：问题陈述"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, Colors.BACKGROUND)

    add_title(slide, "创作之痛")

    # 三大痛点卡片
    pain_points = [
        ("3-6个月", "创作周期"),
        ("10人团队", "人力需求"),
        ("50万元", "制作成本"),
    ]

    card_width = Inches(2.2)
    card_height = Inches(2.5)
    start_y = Inches(2.5)
    start_x = Inches(0.8)
    spacing = Inches(0.3)

    for i, (value, label) in enumerate(pain_points):
        x = start_x + i * (card_width + spacing)

        # 卡片背景
        card = create_gradient_box(slide, x, start_y, card_width, card_height,
                                   Colors.SCREENWRITER, Colors.STORYBOARD)

        # 数值
        textbox = slide.shapes.add_textbox(x, start_y + Inches(0.3),
                                           card_width, Inches(0.8))
        tf = textbox.text_frame
        tf.text = value
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = Colors.TEXT_PRIMARY
        p.font.name = "Arial"

        # 标签
        textbox2 = slide.shapes.add_textbox(x, start_y + Inches(1.3),
                                            card_width, Inches(0.5))
        tf2 = textbox2.text_frame
        tf2.text = label
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        p2.font.size = Pt(16)
        p2.font.color.rgb = Colors.TEXT_PRIMARY
        p2.font.name = "微软雅黑"

    # 解决方案
    add_centered_text(slide, "DramaAI 解决方案", Inches(5.5), font_size=20,
                     color=Colors.ACCENT)
    add_centered_text(slide, "3小时 · 1个人 · 几百元", Inches(5.9), font_size=24,
                     color=Colors.TEXT_PRIMARY)

def create_solution_slide(prs):
    """第3页：解决方案"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, Colors.BACKGROUND)

    add_title(slide, "五阶段AI创作流水线")

    # 中央五边形布局（简化为圆形排列）
    center_x = Inches(5)
    center_y = Inches(3.5)
    radius = Inches(1.3)

    stages = [
        ("编剧", "剧本大纲", Colors.SCREENWRITER),
        ("分镜", "镜头拆解", Colors.STORYBOARD),
        ("设计", "提示词生成", Colors.DESIGNER),
        ("美工", "参考图生成", Colors.ARTIST),
        ("导演", "视频生成", Colors.DIRECTOR),
    ]

    import math

    for i, (title, desc, color) in enumerate(stages):
        angle = math.radians(90 - i * 72)  # 从顶部开始
        x = center_x + Inches(1.8) * math.cos(angle)
        y = center_y + Inches(1.8) * math.sin(angle)

        # 节点圆
        node_size = Inches(1.1)
        node = create_gradient_box(slide, x - node_size/2, y - node_size/2,
                                    node_size, node_size, color, color)

        # 标题
        textbox = slide.shapes.add_textbox(x - node_size/2, y - Inches(0.2),
                                           node_size, Inches(0.4))
        tf = textbox.text_frame
        tf.text = title
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = Colors.BACKGROUND
        p.font.name = "微软雅黑"

        # 描述
        textbox2 = slide.shapes.add_textbox(x - node_size/2, y + Inches(0.15),
                                            node_size, Inches(0.3))
        tf2 = textbox2.text_frame
        tf2.text = desc
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        p2.font.size = Pt(10)
        p2.font.color.rgb = Colors.BACKGROUND
        p2.font.name = "微软雅黑"

    # 中心Logo
    center_box = create_gradient_box(slide, center_x - Inches(0.6), center_y - Inches(0.4),
                                      Inches(1.2), Inches(0.8), Colors.ACCENT, Colors.ACCENT)

    textbox = slide.shapes.add_textbox(center_x - Inches(0.6), center_y - Inches(0.15),
                                       Inches(1.2), Inches(0.4))
    tf = textbox.text_frame
    tf.text = "DramaAI"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = Colors.BACKGROUND
    p.font.name = "Arial"

    # 创新点
    innovations = ["StaleTracker 数据管理", "无限画布 + AI对话", "流式AI响应系统"]
    y_pos = Inches(6.2)
    for innovation in innovations:
        add_centered_text(slide, innovation, y_pos, font_size=14, color=Colors.TEXT_SECONDARY)
        y_pos += Inches(0.35)

def create_tech_slide(prs):
    """第4页：技术架构"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, Colors.BACKGROUND)

    add_title(slide, "技术架构")

    layers = [
        ("前端层", "React 18 + TypeScript + Vite + tldraw", Colors.SCREENWRITER),
        ("状态层", "ProjectContext + IndexedDB", Colors.DESIGNER),
        ("API层", "Vercel Serverless Functions", Colors.ARTIST),
        ("AI层", "Claude 4.5 + RunComfy", Colors.DIRECTOR),
    ]

    box_width = Inches(6)
    box_height = Inches(0.7)
    start_y = Inches(2.2)
    spacing = Inches(0.15)

    for title, content, color in layers:
        # 背景
        box = create_gradient_box(slide, Inches(1.5), start_y, box_width, box_height,
                                   color, color)

        # 标题
        textbox = slide.shapes.add_textbox(Inches(1.6), start_y + Inches(0.1),
                                           Inches(1.5), Inches(0.5))
        tf = textbox.text_frame
        tf.text = title
        p = tf.paragraphs[0]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = Colors.BACKGROUND
        p.font.name = "微软雅黑"

        # 内容
        textbox2 = slide.shapes.add_textbox(Inches(3.2), start_y + Inches(0.15),
                                            Inches(4.5), Inches(0.4))
        tf2 = textbox2.text_frame
        tf2.text = content
        p2 = tf2.paragraphs[0]
        p2.font.size = Pt(12)
        p2.font.color.rgb = Colors.BACKGROUND
        p2.font.name = "Arial"

        start_y += box_height + spacing

    # 技术亮点
    add_centered_text(slide, "全球首个基于无限画布的AI漫剧创作平台",
                     Inches(6), font_size=16, color=Colors.ACCENT)

def create_roi_slide(prs):
    """第5页：ROI分析"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, Colors.BACKGROUND)

    add_title(slide, "价值提升")

    # 三大指标
    metrics = [
        ("720倍", "时间效率", "3个月 → 3小时"),
        ("10倍", "人力效率", "10人 → 1人"),
        ("1000倍+", "成本效率", "50万 → 几百元"),
    ]

    card_width = Inches(2.5)
    card_height = Inches(2)
    start_y = Inches(2.5)
    start_x = Inches(0.7)
    spacing = Inches(0.4)

    for i, (value, label, detail) in enumerate(metrics):
        x = start_x + i * (card_width + spacing)

        card = create_gradient_box(slide, x, start_y, card_width, card_height,
                                   Colors.SCREENWRITER, Colors.STORYBOARD)

        # 数值
        textbox = slide.shapes.add_textbox(x, start_y + Inches(0.2),
                                           card_width, Inches(0.6))
        tf = textbox.text_frame
        tf.text = value
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = Colors.TEXT_PRIMARY
        p.font.name = "Arial"

        # 标签
        textbox2 = slide.shapes.add_textbox(x, start_y + Inches(0.9),
                                            card_width, Inches(0.3))
        tf2 = textbox2.text_frame
        tf2.text = label
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        p2.font.size = Pt(14)
        p2.font.color.rgb = Colors.TEXT_PRIMARY
        p2.font.name = "微软雅黑"

        # 详情
        textbox3 = slide.shapes.add_textbox(x, start_y + Inches(1.3),
                                            card_width, Inches(0.3))
        tf3 = textbox3.text_frame
        tf3.text = detail
        p3 = tf3.paragraphs[0]
        p3.alignment = PP_ALIGN.CENTER
        p3.font.size = Pt(11)
        p3.font.italic = True
        p3.font.color.rgb = Colors.TEXT_PRIMARY
        p3.font.name = "微软雅黑"

    # 商业预测
    add_centered_text(slide, "商业价值预测", Inches(5.2), font_size=18, color=Colors.ACCENT)

    predictions = ["第一年: 60万", "第二年: 600万", "第三年: 6000万"]
    y_pos = Inches(5.7)
    for pred in predictions:
        add_centered_text(slide, pred, y_pos, font_size=14, color=Colors.TEXT_PRIMARY)
        y_pos += Inches(0.3)

def create_roadmap_slide(prs):
    """第6页：发展路线"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, Colors.BACKGROUND)

    add_title(slide, "发展路线图")

    phases = [
        ("短期", "3-6个月", [
            "图片生成功能",
            "视频生成功能",
            "用户系统",
            "云端同步",
        ], Colors.SCREENWRITER),
        ("中期", "6-12个月", [
            "AI模型市场",
            "协作功能",
            "移动端适配",
            "API开放",
        ], Colors.DESIGNER),
        ("长期", "1-3年", [
            "成为AI时代的Adobe",
            "多领域扩展",
            "国际化布局",
            "企业服务",
        ], Colors.ARTIST),
    ]

    card_width = Inches(2.7)
    card_height = Inches(3)
    start_y = Inches(2)
    start_x = Inches(0.5)
    spacing = Inches(0.4)

    for phase_title, time, items, color in phases:
        x = start_x + phases.index((phase_title, time, items, color)) * (card_width + spacing)

        card = create_gradient_box(slide, x, start_y, card_width, card_height,
                                   color, color)

        # 阶段标题
        textbox = slide.shapes.add_textbox(x, start_y + Inches(0.15),
                                           card_width, Inches(0.3))
        tf = textbox.text_frame
        tf.text = phase_title
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = Colors.TEXT_PRIMARY
        p.font.name = "微软雅黑"

        # 时间
        textbox2 = slide.shapes.add_textbox(x, start_y + Inches(0.5),
                                            card_width, Inches(0.25))
        tf2 = textbox2.text_frame
        tf2.text = time
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        p2.font.size = Pt(12)
        p2.font.color.rgb = Colors.TEXT_PRIMARY
        p2.font.name = "微软雅黑"

        # 列表
        item_y = start_y + Inches(0.95)
        for item in items:
            textbox3 = slide.shapes.add_textbox(x + Inches(0.2), item_y,
                                                card_width - Inches(0.4), Inches(0.25))
            tf3 = textbox3.text_frame
            tf3.text = f"• {item}"
            p3 = tf3.paragraphs[0]
            p3.font.size = Pt(10)
            p3.font.color.rgb = Colors.TEXT_PRIMARY
            p3.font.name = "微软雅黑"
            item_y += Inches(0.35)

    add_centered_text(slide, "愿景：让每个人都能成为动漫导演",
                     Inches(5.8), font_size=14, color=Colors.ACCENT)

def create_contact_slide(prs):
    """第7页：联系页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, Colors.BACKGROUND)

    # 大标题
    textbox = slide.shapes.add_textbox(Inches(1.5), Inches(2.5),
                                       Inches(7), Inches(1))
    tf = textbox.text_frame
    tf.text = "DramaAI"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = Colors.TEXT_PRIMARY
    p.font.name = "Arial"

    # Slogan
    add_centered_text(slide, "让每个人都能成为动漫导演", Inches(3.8),
                     font_size=24, color=Colors.ACCENT)

    # 装饰圈
    shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(3.5), Inches(2.2), Inches(3), Inches(2)
    )
    shape.fill.background()
    shape.line.color.rgb = Colors.SCREENWRITER
    shape.line.width = Pt(3)

    # 联系方式
    contact_info = [
        ("网站", "https://animeaistudio.com"),
        ("邮箱", "contact@animeaistudio.com"),
        ("GitHub", "github.com/dramaai/dramaai"),
    ]

    y_start = Inches(4.8)
    for label, value in contact_info:
        textbox = slide.shapes.add_textbox(Inches(2), y_start,
                                           Inches(6), Inches(0.3))
        tf = textbox.text_frame
        tf.text = f"{label}: {value}"
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(14)
        p.font.color.rgb = Colors.TEXT_SECONDARY
        p.font.name = "微软雅黑"
        y_start += Inches(0.35)

    # 感谢
    add_centered_text(slide, "感谢聆听", Inches(6.3), font_size=20, color=Colors.TEXT_PRIMARY)
    add_centered_text(slide, "Q & A", Inches(6.6), font_size=16, color=Colors.TEXT_SECONDARY)

def create_pitch_deck():
    """创建完整的PPT"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 生成各页
    create_cover_slide(prs)
    create_problem_slide(prs)
    create_solution_slide(prs)
    create_tech_slide(prs)
    create_roi_slide(prs)
    create_roadmap_slide(prs)
    create_contact_slide(prs)

    # 保存
    prs.save(OUTPUT_FILE)

    print("="*60)
    print("[SUCCESS] 已生成世界级水准PPT!")
    print("="*60)
    print(f"文件: {OUTPUT_FILE}")
    print(f"页数: 7页")
    print(f"尺寸: 10 x 7.5 英寸")
    print(f"设计理念: Data Flow Chromatics")
    print("="*60)
    print("\n页面列表:")
    print("  1. 封面 - DramaAI")
    print("  2. 创作之痛 - 问题陈述")
    print("  3. 五阶段AI创作流水线")
    print("  4. 技术架构")
    print("  5. 价值提升 (ROI)")
    print("  6. 发展路线图")
    print("  7. 联系方式")
    print("="*60)

if __name__ == "__main__":
    create_pitch_deck()
